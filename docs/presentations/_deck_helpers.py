"""Shared slide-building helpers for Quidnug presentation decks.

Each deck's build_deck.py imports from this module via:

    import sys, pathlib
    sys.path.insert(0, str(pathlib.Path(__file__).parent.parent))
    from _deck_helpers import *   # noqa

Palette, fonts, slide-template functions, and the defRPr
color-override workaround live here so the three decks render
consistently and any improvement to the helpers propagates to all
of them.
"""
import pathlib
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# ---- Palette -----------------------------------------------------------------
DARK_BG = RGBColor(0x0A, 0x16, 0x28)
DARK_CARD = RGBColor(0x14, 0x21, 0x37)
MID = RGBColor(0x1C, 0x3A, 0x5E)
TEAL = RGBColor(0x00, 0xD4, 0xA8)
TEAL_SOFT = RGBColor(0xC3, 0xEF, 0xE3)
CORAL = RGBColor(0xFF, 0x46, 0x55)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
TEXT_DARK = RGBColor(0x1A, 0x1D, 0x23)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SW, SH = Inches(13.333), Inches(7.5)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"


def _color_hex(rgb_color):
    return str(rgb_color).upper()


def _force_color_override(element, rgb_hex):
    """Inject defRPr with srgbClr into a paragraph pPr so editors
    (e.g. OpenOffice Impress) that prefer paragraph-default colors
    over run-level colors still render the intended color."""
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    A = f"{{{ns}}}"
    p_pr = element.find(f"{A}pPr")
    if p_pr is None:
        p_pr = etree.SubElement(element, f"{A}pPr")
        element.insert(0, p_pr)
    for old in p_pr.findall(f"{A}defRPr"):
        p_pr.remove(old)
    def_r = etree.SubElement(p_pr, f"{A}defRPr")
    sf = etree.SubElement(def_r, f"{A}solidFill")
    rgb = etree.SubElement(sf, f"{A}srgbClr")
    rgb.set("val", rgb_hex)


def set_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_keyline(slide, color=TEAL, width=Inches(0.1)):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width, SH)
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = color


def add_footer(slide, slide_num, total, brand, text_color=TEXT_MUTED):
    box = slide.shapes.add_textbox(Inches(12.3), Inches(7.1),
                                   Inches(1.0), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{slide_num} / {total}"
    run.font.size = Pt(9)
    run.font.name = BODY_FONT
    run.font.color.rgb = text_color
    _force_color_override(p._p, _color_hex(text_color))

    box2 = slide.shapes.add_textbox(Inches(0.3), Inches(7.1),
                                    Inches(6.5), Inches(0.3))
    tf = box2.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = brand
    run.font.size = Pt(9)
    run.font.name = BODY_FONT
    run.font.color.rgb = text_color
    _force_color_override(p._p, _color_hex(text_color))


def add_title(slide, text, x=Inches(0.6), y=Inches(0.5),
              w=Inches(12.1), h=Inches(1.0),
              size=36, color=TEXT_DARK, bold=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = HEADER_FONT
    run.font.color.rgb = color
    _force_color_override(p._p, _color_hex(color))
    return box


def add_subtitle(slide, text, x=Inches(0.6), y=Inches(1.45),
                 w=Inches(12.1), h=Inches(0.6),
                 size=18, color=TEXT_MUTED, italic=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.name = BODY_FONT
    run.font.color.rgb = color
    _force_color_override(p._p, _color_hex(color))
    return box


def add_bullets(slide, items, x=Inches(0.6), y=Inches(2.2),
                w=Inches(12.1), h=Inches(4.5),
                size=16, color=TEXT_DARK, bullet_color=TEAL,
                line_gap=6):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(line_gap)

        bullet_run = p.add_run()
        bullet_run.text = "\u25B8 "
        bullet_run.font.size = Pt(size)
        bullet_run.font.name = BODY_FONT
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.bold = True

        if isinstance(item, tuple):
            strong, rest = item
            r1 = p.add_run()
            r1.text = strong
            r1.font.size = Pt(size)
            r1.font.name = BODY_FONT
            r1.font.color.rgb = color
            r1.font.bold = True

            r2 = p.add_run()
            r2.text = " " + rest if rest else ""
            r2.font.size = Pt(size)
            r2.font.name = BODY_FONT
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.name = BODY_FONT
            r.font.color.rgb = color
        _force_color_override(p._p, _color_hex(color))

    return box


def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    if isinstance(notes_text, str):
        notes_text = [notes_text]
    for i, line in enumerate(notes_text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(12)
        run.font.name = BODY_FONT


def add_card(slide, x, y, w, h, fill=CARD_BG, border=CARD_BORDER,
             border_width=Pt(1.0), corner=True):
    if corner:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.adjustments[0] = 0.06
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = border_width
    return shape


def add_text_in(slide, x, y, w, h, text, size=14,
                color=TEXT_DARK, bold=False, align="left",
                italic=False, font=BODY_FONT, anchor="top"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    _force_color_override(p._p, _color_hex(color))
    return box


def add_image(slide, path, x, y, w=None, h=None):
    kwargs = {}
    if w is not None:
        kwargs["width"] = w
    if h is not None:
        kwargs["height"] = h
    return slide.shapes.add_picture(str(path), x, y, **kwargs)


def add_table(slide, data, x, y, w, h,
              header_fill=TEAL, header_fg=WHITE,
              row_fill=WHITE, alt_fill=LIGHT_BG,
              border=CARD_BORDER, header_size=13, body_size=12):
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table

    for c, header in enumerate(data[0]):
        cell = tbl.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = str(header)
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.name = BODY_FONT
        run.font.color.rgb = header_fg
        _force_color_override(p._p, _color_hex(header_fg))

    for r in range(1, rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = alt_fill if r % 2 == 0 else row_fill
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(data[r][c])
            run.font.size = Pt(body_size)
            run.font.name = BODY_FONT
            run.font.color.rgb = TEXT_DARK
            _force_color_override(p._p, _color_hex(TEXT_DARK))

    return tbl


def add_quote(slide, text, attribution, x=Inches(1.5), y=Inches(2.5),
              w=Inches(10.3), size=22, color=TEXT_DARK):
    mark = slide.shapes.add_textbox(x, y - Inches(0.4), Inches(1),
                                    Inches(1))
    tf = mark.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "\u201C"
    r.font.size = Pt(96)
    r.font.name = HEADER_FONT
    r.font.color.rgb = TEAL
    _force_color_override(p._p, _color_hex(TEAL))

    box = slide.shapes.add_textbox(x + Inches(0.7), y, w, Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.name = HEADER_FONT
    run.font.color.rgb = color
    _force_color_override(p._p, _color_hex(color))

    p2 = tf.add_paragraph()
    p2.space_before = Pt(16)
    r2 = p2.add_run()
    r2.text = "\u2014 " + attribution
    r2.font.size = Pt(14)
    r2.font.name = BODY_FONT
    r2.font.color.rgb = TEXT_MUTED
    _force_color_override(p2._p, _color_hex(TEXT_MUTED))


def add_big_stat(slide, stat, label, context=None,
                 y=Inches(2.5), stat_color=TEAL,
                 label_color=TEXT_DARK):
    box = slide.shapes.add_textbox(Inches(0.6), y, Inches(12.1),
                                   Inches(2.8))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = stat
    run.font.size = Pt(120)
    run.font.bold = True
    run.font.name = HEADER_FONT
    run.font.color.rgb = stat_color
    _force_color_override(p._p, _color_hex(stat_color))

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)
    run = p2.add_run()
    run.text = label
    run.font.size = Pt(24)
    run.font.name = BODY_FONT
    run.font.color.rgb = label_color
    _force_color_override(p2._p, _color_hex(label_color))

    if context:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(12)
        run = p3.add_run()
        run.text = context
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.name = BODY_FONT
        run.font.color.rgb = TEXT_MUTED
        _force_color_override(p3._p, _color_hex(TEXT_MUTED))


def make_presentation():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


# ---- Slide templates ---------------------------------------------------------
def title_slide(prs, title, subtitle, eyebrow=None, deck_brand=""):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5),
                                  Inches(0.6), Inches(2.5))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL

    for i, (cx, cy, sz) in enumerate([(11.5, 5.8, 0.3), (12.1, 6.0, 0.2),
                                      (12.6, 5.6, 0.15), (12.5, 6.4, 0.12)]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx), Inches(cy), Inches(sz), Inches(sz))
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL if i % 2 == 0 else TEAL_SOFT

    if eyebrow:
        add_text_in(slide, Inches(1.0), Inches(1.5), Inches(10),
                    Inches(0.4),
                    eyebrow, size=13, color=TEAL, bold=True,
                    font=BODY_FONT)

    add_title(slide, title,
              x=Inches(1.0), y=Inches(2.3), w=Inches(11.5),
              h=Inches(2.0),
              size=48, color=WHITE, bold=True)

    add_subtitle(slide, subtitle,
                 x=Inches(1.0), y=Inches(4.8), w=Inches(11.5),
                 h=Inches(1.6),
                 size=20, color=TEAL_SOFT, italic=False)

    return slide


def section_divider(prs, section_num, title, subtitle=None):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.3),
                                  SW, Inches(0.08))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL

    add_text_in(slide, Inches(0.8), Inches(2.0), Inches(4), Inches(1),
                f"Section {section_num}", size=18, color=TEAL, bold=True)

    add_title(slide, title,
              x=Inches(0.8), y=Inches(2.7), w=Inches(11.7),
              h=Inches(2.0),
              size=44, color=WHITE)

    if subtitle:
        add_subtitle(slide, subtitle,
                     x=Inches(0.8), y=Inches(4.8), w=Inches(11.7),
                     h=Inches(1),
                     size=20, color=TEAL_SOFT, italic=False)

    return slide


def content_slide(prs, title, bullets=None, subtitle=None, notes=None,
                  image=None, image_pos=None, assets_dir=None):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    if image and assets_dir:
        if image_pos is None:
            image_pos = (Inches(0.6), Inches(2.3),
                         Inches(12.1), None)
        x, y, w, h = image_pos
        add_image(slide, assets_dir / image, x, y, w=w, h=h)
    if bullets:
        add_bullets(slide, bullets)

    if notes:
        add_notes(slide, notes)
    return slide


def two_col_slide(prs, title, left_title, left_items,
                  right_title, right_items, subtitle=None, notes=None,
                  left_color=EMERALD, right_color=CORAL):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    add_card(slide, Inches(0.6), Inches(2.3),
             Inches(6.0), Inches(4.6))
    lt = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.3),
        Inches(6.0), Inches(0.55))
    lt.line.fill.background()
    lt.fill.solid()
    lt.fill.fore_color.rgb = left_color
    add_text_in(slide, Inches(0.6), Inches(2.3), Inches(6.0),
                Inches(0.55),
                left_title, size=15, color=WHITE, bold=True,
                align="center", anchor="middle")

    lbox = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), Inches(5.6), Inches(3.8))
    tf = lbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        b = p.add_run()
        b.text = "\u2022 "
        b.font.size = Pt(13)
        b.font.bold = True
        b.font.color.rgb = left_color
        r = p.add_run()
        r.text = item
        r.font.size = Pt(13)
        r.font.name = BODY_FONT
        r.font.color.rgb = TEXT_DARK
        _force_color_override(p._p, _color_hex(TEXT_DARK))

    add_card(slide, Inches(6.75), Inches(2.3),
             Inches(6.0), Inches(4.6))
    rt = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.75), Inches(2.3),
        Inches(6.0), Inches(0.55))
    rt.line.fill.background()
    rt.fill.solid()
    rt.fill.fore_color.rgb = right_color
    add_text_in(slide, Inches(6.75), Inches(2.3), Inches(6.0),
                Inches(0.55), right_title, size=15, color=WHITE,
                bold=True, align="center", anchor="middle")
    rbox = slide.shapes.add_textbox(
        Inches(6.95), Inches(3.0), Inches(5.6), Inches(3.8))
    tf = rbox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        b = p.add_run()
        b.text = "\u2022 "
        b.font.size = Pt(13)
        b.font.bold = True
        b.font.color.rgb = right_color
        r = p.add_run()
        r.text = item
        r.font.size = Pt(13)
        r.font.name = BODY_FONT
        r.font.color.rgb = TEXT_DARK
        _force_color_override(p._p, _color_hex(TEXT_DARK))

    if notes:
        add_notes(slide, notes)
    return slide


def stat_slide(prs, stat, label, context=None, subtitle=None,
               title=None, notes=None, stat_color=TEAL):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    if title:
        add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)
    add_big_stat(slide, stat, label, context=context,
                 y=Inches(2.4) if title else Inches(1.6),
                 stat_color=stat_color)
    if notes:
        add_notes(slide, notes)
    return slide


def quote_slide(prs, quote, attrib, context=None, notes=None,
                title=None):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    if title:
        add_title(slide, title)
    add_quote(slide, quote, attrib,
              x=Inches(1.2), y=Inches(2.6),
              w=Inches(10.8), size=24)
    if context:
        add_text_in(slide, Inches(1.2), Inches(6.3), Inches(10.8),
                    Inches(0.6), context,
                    size=12, color=TEXT_MUTED, italic=True)
    if notes:
        add_notes(slide, notes)
    return slide


def table_slide(prs, title, data, subtitle=None, notes=None,
                col_widths=None, header_size=13, body_size=12):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    rows = len(data)
    table_h = Inches(min(4.5, 0.5 + rows * 0.45))
    tbl = add_table(slide, data,
                    Inches(0.6), Inches(2.3),
                    Inches(12.1), table_h,
                    header_size=header_size, body_size=body_size)
    if col_widths:
        total = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(12.1 * w / total)

    if notes:
        add_notes(slide, notes)
    return slide


def image_slide(prs, title, image, caption=None, subtitle=None,
                notes=None, image_y=Inches(2.0),
                image_h=Inches(4.6), assets_dir=None):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    pic = slide.shapes.add_picture(
        str(assets_dir / image), Inches(0.8),
        image_y, height=image_h)
    pic.left = int((SW - pic.width) / 2)

    if caption:
        add_text_in(slide, Inches(0.6), Inches(6.8), Inches(12.1),
                    Inches(0.3), caption,
                    size=10, color=TEXT_MUTED, italic=True,
                    align="center")

    if notes:
        add_notes(slide, notes)
    return slide


def code_slide(prs, title, code_lines, subtitle=None, notes=None,
               caption=None):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6), Inches(2.3),
        Inches(12.1), Inches(4.3))
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_BG
    card.line.color.rgb = MID
    card.line.width = Pt(0.8)

    box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.5), Inches(11.5), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        if line.strip().startswith("//") or line.strip().startswith("#"):
            color = RGBColor(0x94, 0xA3, 0xB8)
        elif line.strip().startswith('"') and ":" in line:
            color = TEAL
        else:
            color = WHITE
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.name = "Consolas"
        r.font.color.rgb = color
        _force_color_override(p._p, _color_hex(color))

    if caption:
        add_text_in(slide, Inches(0.6), Inches(6.8), Inches(12.1),
                    Inches(0.3), caption,
                    size=10, color=TEXT_MUTED, italic=True,
                    align="center")

    if notes:
        add_notes(slide, notes)
    return slide


def icon_grid_slide(prs, title, items, subtitle=None, notes=None,
                    cols=3):
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_BG)
    add_keyline(slide)
    add_title(slide, title)
    if subtitle:
        add_subtitle(slide, subtitle)

    rows = (len(items) + cols - 1) // cols
    card_w = 12.1 / cols - 0.15
    card_h = 4.4 / rows - 0.15

    for idx, item in enumerate(items):
        if len(item) == 3:
            label, desc, color = item
        else:
            label, desc = item
            color = TEAL
        r = idx // cols
        c = idx % cols
        x = 0.6 + c * (card_w + 0.15)
        y = 2.3 + r * (card_h + 0.15)

        add_card(slide, Inches(x), Inches(y),
                 Inches(card_w), Inches(card_h))
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
            Inches(card_w), Inches(0.18))
        strip.line.fill.background()
        strip.fill.solid()
        strip.fill.fore_color.rgb = color

        add_text_in(slide, Inches(x + 0.2), Inches(y + 0.3),
                    Inches(card_w - 0.4), Inches(0.5),
                    label, size=15, color=TEXT_DARK, bold=True)
        add_text_in(slide, Inches(x + 0.2), Inches(y + 0.85),
                    Inches(card_w - 0.4), Inches(card_h - 0.95),
                    desc, size=11.5, color=TEXT_MUTED)

    if notes:
        add_notes(slide, notes)
    return slide


def closing_slide(prs, title, subtitle=None, cta=None,
                  resources=None, notes=None):
    slide = blank_slide(prs)
    set_bg(slide, DARK_BG)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                  Inches(0.5), SH)
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = TEAL

    add_title(slide, title, x=Inches(1.0), y=Inches(0.6),
              w=Inches(11.5), h=Inches(1.0),
              size=40, color=WHITE)
    if subtitle:
        add_subtitle(slide, subtitle, x=Inches(1.0), y=Inches(1.7),
                     w=Inches(11.5), h=Inches(0.5),
                     size=18, color=TEAL_SOFT, italic=False)

    if cta:
        add_text_in(slide, Inches(1.0), Inches(2.5), Inches(11.5),
                    Inches(2.0), cta,
                    size=16, color=WHITE, bold=False)

    if resources:
        add_text_in(slide, Inches(1.0), Inches(4.7),
                    Inches(11.5), Inches(0.35),
                    "Resources", size=12, color=TEAL, bold=True)
        y = 5.05
        for line in resources:
            add_text_in(slide, Inches(1.0), Inches(y),
                        Inches(11.5), Inches(0.3),
                        line, size=12, color=TEAL_SOFT)
            y += 0.29

    if notes:
        add_notes(slide, notes)
    return slide
