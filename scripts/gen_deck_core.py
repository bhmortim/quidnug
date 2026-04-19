"""Design system + reusable slide primitive helpers for the Quidnug deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree

# --- Slide canvas (16:9 widescreen) ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# --- Palette: cryptographic-trust themed (not generic blue) ---
MIDNIGHT   = RGBColor(0x0B, 0x19, 0x29)   # Deep midnight — dark bg / text
NAVY       = RGBColor(0x1E, 0x3A, 0x5F)   # Primary (~60%)
TEAL       = RGBColor(0x14, 0xB8, 0xA6)   # Accent (~20%)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)   # Sharp contrast (~10%)
ICE        = RGBColor(0xF1, 0xF5, 0xF9)   # Soft light fill
CLOUD      = RGBColor(0xE2, 0xE8, 0xF0)   # A touch darker than ICE
SLATE      = RGBColor(0x64, 0x74, 0x8B)   # Muted text / borders
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
RED        = RGBColor(0xEF, 0x44, 0x44)
SOFT_TEAL  = RGBColor(0xCC, 0xFB, 0xF1)
SOFT_AMBER = RGBColor(0xFE, 0xF3, 0xC7)

# --- Typography ---
H_FONT = "Georgia"       # Classical, trustworthy
B_FONT = "Calibri"       # Clean, widely available
CODE_FONT = "Consolas"

# --- Helpers: primitives ---

def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, font=B_FONT, size=14, color=MIDNIGHT,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)

    lines = text if isinstance(text, list) else [text]
    for i, raw in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = raw
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=MIDNIGHT, bullet_char="•",
            font=B_FONT, indent_unit=0.22):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)

    for i, item in enumerate(items):
        # item can be str or (text, level=0, bold=False, accent=False)
        if isinstance(item, tuple):
            txt = item[0]
            level = item[1] if len(item) > 1 else 0
            bold = item[2] if len(item) > 2 else False
            accent = item[3] if len(item) > 3 else False
        else:
            txt = item
            level = 0
            bold = False
            accent = False

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0  # we handle indent manually
        p.space_after = Pt(4)

        indent = " " * (level * 3)
        prefix = "›" if level == 1 else bullet_char
        run = p.add_run()
        run.text = f"{indent}{prefix}  {txt}"
        run.font.name = font
        run.font.size = Pt(size - (1 if level else 0))
        run.font.color.rgb = (TEAL if accent else color)
        run.font.bold = bold
    return tb


def rect(slide, x, y, w, h, fill, line=None, line_w=0.75, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    if not shadow:
        shp.shadow.inherit = False
    return shp


def rounded(slide, x, y, w, h, fill, line=None, line_w=0.75, corner=0.15, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    try:
        shp.adjustments[0] = corner
    except Exception:
        pass
    if not shadow:
        shp.shadow.inherit = False
    return shp


def hexagon(slide, x, y, w, h, fill, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def oval(slide, x, y, w, h, fill, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def shape_text(shp, text, font=B_FONT, size=14, color=WHITE, bold=False,
               italic=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)

    lines = text if isinstance(text, list) else [text]
    tf.paragraphs[0].text = ""
    for i, raw in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = raw
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic


def arrow(slide, x1, y1, x2, y2, color=SLATE, width=1.5, end_arrow=True):
    """Straight arrow from (x1,y1) to (x2,y2). All EMU/Inches."""
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln = conn.line
    ln.color.rgb = color
    ln.width = Pt(width)
    if end_arrow:
        # Add arrowhead via XML
        lnXml = conn.line._get_or_add_ln()
        # head
        head = etree.SubElement(
            lnXml,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd"
        )
        head.set("type", "none")
        tail = etree.SubElement(
            lnXml,
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd"
        )
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("h", "med")
    return conn


def notes(slide, text):
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.text = text


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def hex_watermark(slide, color=TEAL):
    """Subtle hex mark in upper-right corner — our motif."""
    h = hexagon(slide, Inches(12.55), Inches(0.25), Inches(0.5), Inches(0.42),
                fill=None, line=color, line_w=1.25)
    return h


def page_chrome(slide, title, kicker=None, on_dark=False, page_num=None, total=None):
    """Title bar, kicker, hex watermark, thin divider, footer."""
    title_color = WHITE if on_dark else MIDNIGHT
    kicker_color = TEAL if on_dark else TEAL

    if kicker:
        text_box(slide, Inches(0.5), Inches(0.35), Inches(10), Inches(0.35),
                 kicker, font=B_FONT, size=12, color=kicker_color, bold=True)
        text_box(slide, Inches(0.5), Inches(0.7), Inches(12), Inches(0.7),
                 title, font=H_FONT, size=30, color=title_color, bold=True)
    else:
        text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8),
                 title, font=H_FONT, size=32, color=title_color, bold=True)

    hex_watermark(slide, color=(TEAL if on_dark else TEAL))

    # Footer
    footer_color = CLOUD if on_dark else SLATE
    text_box(slide, Inches(0.5), Inches(7.15), Inches(10), Inches(0.3),
             "Quidnug — A Decentralized Protocol for Relational Trust",
             font=B_FONT, size=9, color=footer_color, italic=True)
    if page_num is not None:
        label = f"{page_num}" if total is None else f"{page_num} / {total}"
        text_box(slide, Inches(11.8), Inches(7.15), Inches(1.3), Inches(0.3),
                 label, font=B_FONT, size=9, color=footer_color, align=PP_ALIGN.RIGHT)


# --- Higher-level slide builders ---

def slide_title(prs):
    s = blank(prs)
    set_bg(s, MIDNIGHT)
    # Decorative hexes in upper right
    hexagon(s, Inches(11.0), Inches(0.4), Inches(0.7), Inches(0.6), NAVY, line=TEAL, line_w=1)
    hexagon(s, Inches(11.7), Inches(0.85), Inches(0.45), Inches(0.4), TEAL)
    hexagon(s, Inches(12.4), Inches(0.4), Inches(0.55), Inches(0.48), None, line=TEAL, line_w=1)
    # Decorative hex tile in left-lower area
    hexagon(s, Inches(0.6), Inches(6.4), Inches(0.45), Inches(0.4), AMBER)
    hexagon(s, Inches(1.2), Inches(6.25), Inches(0.65), Inches(0.55), None, line=AMBER, line_w=1)

    text_box(s, Inches(0), Inches(2.2), Inches(13.333), Inches(1.4),
             "Quidnug", font=H_FONT, size=110, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
    text_box(s, Inches(0), Inches(3.85), Inches(13.333), Inches(0.6),
             "A Decentralized Protocol for Relational Trust",
             font=H_FONT, size=28, color=TEAL, italic=True, align=PP_ALIGN.CENTER)
    rect(s, Inches(5.66), Inches(4.7), Inches(2), Inches(0.035), TEAL)
    text_box(s, Inches(0), Inches(4.95), Inches(13.333), Inches(0.45),
             "Identity  ·  Ownership  ·  Auditable State",
             font=B_FONT, size=17, color=ICE, align=PP_ALIGN.CENTER)
    text_box(s, Inches(0), Inches(5.55), Inches(13.333), Inches(0.4),
             "Developer Overview  ·  v2026.04",
             font=B_FONT, size=13, color=SLATE, italic=True, align=PP_ALIGN.CENTER)
    return s


def slide_section(prs, number, label, title):
    s = blank(prs)
    set_bg(s, MIDNIGHT)
    # Big glyphs — number + hex
    text_box(s, Inches(0.6), Inches(1.8), Inches(4.5), Inches(4),
             f"{number:02d}", font=H_FONT, size=220, color=TEAL, bold=True,
             align=PP_ALIGN.LEFT)
    # Accent hexes
    hexagon(s, Inches(4.6), Inches(2.6), Inches(0.6), Inches(0.52), AMBER)
    hexagon(s, Inches(5.1), Inches(3.2), Inches(0.45), Inches(0.4), None, line=AMBER, line_w=1.25)

    text_box(s, Inches(5.7), Inches(2.85), Inches(7.3), Inches(0.5),
             label, font=B_FONT, size=16, color=TEAL, bold=True)
    text_box(s, Inches(5.7), Inches(3.2), Inches(7.3), Inches(3),
             title, font=H_FONT, size=44, color=WHITE, bold=True)
    # Bottom accent
    rect(s, Inches(5.7), Inches(5.4), Inches(3), Inches(0.035), TEAL)
    return s


def slide_content(prs, title, kicker=None, page=None, total=None):
    s = blank(prs)
    set_bg(s, WHITE)
    page_chrome(s, title, kicker=kicker, on_dark=False, page_num=page, total=total)
    return s


def slide_dark(prs, title, kicker=None, page=None, total=None):
    s = blank(prs)
    set_bg(s, MIDNIGHT)
    page_chrome(s, title, kicker=kicker, on_dark=True, page_num=page, total=total)
    return s


def slide_quote(prs, quote, attribution=None, page=None, total=None):
    s = blank(prs)
    set_bg(s, NAVY)
    # Big opening hex
    hexagon(s, Inches(0.7), Inches(0.7), Inches(0.7), Inches(0.6), TEAL)
    text_box(s, Inches(0.7), Inches(0.8), Inches(0.7), Inches(0.6),
             "\u201C", font=H_FONT, size=48, color=MIDNIGHT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Quote text
    text_box(s, Inches(1.5), Inches(2.2), Inches(10.5), Inches(3.6),
             quote, font=H_FONT, size=36, color=WHITE, italic=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15)
    if attribution:
        rect(s, Inches(1.5), Inches(6.0), Inches(0.6), Inches(0.03), TEAL)
        text_box(s, Inches(1.5), Inches(6.1), Inches(9), Inches(0.4),
                 f"— {attribution}", font=B_FONT, size=14, color=TEAL, italic=True)

    hex_watermark(s, color=TEAL)
    text_box(s, Inches(0.5), Inches(7.15), Inches(10), Inches(0.3),
             "Quidnug — A Decentralized Protocol for Relational Trust",
             font=B_FONT, size=9, color=CLOUD, italic=True)
    if page is not None:
        label = f"{page}" if total is None else f"{page} / {total}"
        text_box(s, Inches(11.8), Inches(7.15), Inches(1.3), Inches(0.3),
                 label, font=B_FONT, size=9, color=CLOUD, align=PP_ALIGN.RIGHT)
    return s


def slide_bigstat(prs, stat, label, caption=None, page=None, total=None):
    s = blank(prs)
    set_bg(s, MIDNIGHT)
    # Left big stat
    text_box(s, Inches(0.8), Inches(1.8), Inches(7), Inches(3.4),
             stat, font=H_FONT, size=180, color=TEAL, bold=True,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    text_box(s, Inches(0.8), Inches(5.1), Inches(7), Inches(0.6),
             label, font=H_FONT, size=28, color=WHITE, bold=True)
    if caption:
        text_box(s, Inches(0.8), Inches(5.75), Inches(7), Inches(1.3),
                 caption, font=B_FONT, size=16, color=CLOUD, line_spacing=1.25)
    # Decorative cluster on right
    hexagon(s, Inches(10.2), Inches(2.4), Inches(1.3), Inches(1.1), NAVY, line=TEAL, line_w=1.5)
    hexagon(s, Inches(11.4), Inches(3.1), Inches(1.0), Inches(0.85), TEAL)
    hexagon(s, Inches(9.7), Inches(3.6), Inches(0.8), Inches(0.7), None, line=AMBER, line_w=1.25)
    hexagon(s, Inches(10.9), Inches(4.3), Inches(1.1), Inches(0.95), None, line=TEAL, line_w=1.5)

    hex_watermark(s, color=TEAL)
    text_box(s, Inches(0.5), Inches(7.15), Inches(10), Inches(0.3),
             "Quidnug — A Decentralized Protocol for Relational Trust",
             font=B_FONT, size=9, color=CLOUD, italic=True)
    if page is not None:
        label2 = f"{page}" if total is None else f"{page} / {total}"
        text_box(s, Inches(11.8), Inches(7.15), Inches(1.3), Inches(0.3),
                 label2, font=B_FONT, size=9, color=CLOUD, align=PP_ALIGN.RIGHT)
    return s


# --- Composite shape helpers for diagrams ---

def card(slide, x, y, w, h, title, body=None, accent=TEAL, fill=ICE,
         title_color=None, body_color=None, title_size=17, body_size=12.5,
         corner=0.08, badge_text=None):
    """Info card: rounded box with colored accent stripe and optional badge."""
    rounded(slide, x, y, w, h, fill=fill, line=None, corner=corner)
    # Accent stripe on left
    stripe = rect(slide, x, y, Inches(0.09), h, accent)
    # Optional badge
    tc = title_color or MIDNIGHT
    bc = body_color or SLATE
    if badge_text:
        badge_w = Inches(0.85)
        oval(slide, x + w - badge_w - Inches(0.15), y + Inches(0.15),
             badge_w, Inches(0.38), accent)
        text_box(slide, x + w - badge_w - Inches(0.15), y + Inches(0.15),
                 badge_w, Inches(0.38), badge_text, font=B_FONT, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    text_box(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.4),
             Inches(0.45), title, font=H_FONT, size=title_size,
             color=tc, bold=True)
    if body:
        text_box(slide, x + Inches(0.25), y + Inches(0.6), w - Inches(0.4),
                 h - Inches(0.7), body, font=B_FONT, size=body_size,
                 color=bc, line_spacing=1.2)


def chip(slide, x, y, w, h, text, fill=TEAL, color=WHITE, size=11, bold=True):
    rounded(slide, x, y, w, h, fill=fill, line=None, corner=0.5)
    text_box(slide, x, y, w, h, text, font=B_FONT, size=size, color=color,
             bold=bold, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def numbered_circle(slide, x, y, d, num, fill=TEAL, color=WHITE, size=18):
    oval(slide, x, y, Inches(d), Inches(d), fill)
    text_box(slide, x, y, Inches(d), Inches(d), str(num), font=H_FONT,
             size=size, color=color, bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def table_rows(slide, x, y, col_widths, rows, header_fill=NAVY,
               header_color=WHITE, row_alt=ICE, border=CLOUD,
               header_size=12, body_size=11, row_h=Inches(0.42),
               header_h=Inches(0.45), bold_first_col=False):
    """Custom-styled table built from rectangles + textboxes (not native table)."""
    total_w = sum(col_widths)
    # Header
    cur_x = x
    for ci, cw in enumerate(col_widths):
        rect(slide, cur_x, y, cw, header_h, header_fill, line=border, line_w=0.5)
        header_text = rows[0][ci] if rows else ""
        text_box(slide, cur_x, y, cw, header_h, header_text, font=B_FONT,
                 size=header_size, color=header_color, bold=True,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
        cur_x += cw
    # Body
    for ri, row in enumerate(rows[1:]):
        cur_x = x
        ry = y + header_h + row_h * ri
        fill = row_alt if (ri % 2 == 0) else WHITE
        for ci, cw in enumerate(col_widths):
            rect(slide, cur_x, ry, cw, row_h, fill, line=border, line_w=0.5)
            cell = row[ci] if ci < len(row) else ""
            bold = bold_first_col and ci == 0
            text_box(slide, cur_x, ry, cw, row_h, cell, font=B_FONT,
                     size=body_size, color=MIDNIGHT, bold=bold,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            cur_x += cw
    return total_w, header_h + row_h * (len(rows) - 1)


if __name__ == "__main__":
    print("gen_deck_core.py — import, don't run directly.")
